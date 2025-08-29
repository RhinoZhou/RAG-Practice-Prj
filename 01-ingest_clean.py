#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
临床文档数据清洗程序

功能：读取原始数据（EMR、指南、药典、CSV/DB），完成去标识化、去重、格式统一、单位标准化（全半角/大小写/时区），输出干净语料。
输入：原始文本/数据库连接
输出：clean_text.jsonl（含文档ID、段落、元数据）
依赖：pandas, pyarrow, pydantic
"""
import os
import json
import re
import uuid
from datetime import datetime
import pandas as pd
from pydantic import BaseModel, Field
import pdfplumber
from docx import Document
from pptx import Presentation
import hashlib
import logging
from typing import List, Dict, Any, Optional

# 配置日志，同时输出到控制台和文件
# 设置日志级别为INFO，格式包含时间戳、日志级别和消息内容
# 日志将同时保存到文件ingest_clean.log和输出到控制台
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("ingest_clean.log"),  # 文件日志处理器
        logging.StreamHandler()  # 控制台日志处理器
    ]
)
logger = logging.getLogger(__name__)  # 获取日志记录器实例
class DocumentSegment(BaseModel):
    """文档片段模型，用于结构化存储清洗后的文本片段和相关元数据
    
    该模型使用Pydantic进行数据验证和序列化，确保输出数据的一致性和完整性。
    所有字段都有明确的类型定义和描述，便于后续处理和分析。
    """
    doc_id: str = Field(description="文档唯一标识符，使用UUID生成，确保全局唯一性")
    segment_id: str = Field(description="段落唯一标识符，使用UUID生成，确保全局唯一性")
    text: str = Field(description="清洗后的文本内容，经过去标识化、格式统一等处理")
    original_path: str = Field(description="原始文件路径，记录数据来源")
    file_type: str = Field(description="文件类型，如pdf、docx、pptx等")
    page_num: Optional[int] = Field(default=None, description="页码，适用于PDF和PPTX文件")
    section: Optional[str] = Field(default=None, description="章节信息，根据文件类型不同有不同表示")
    created_at: str = Field(default_factory=lambda: datetime.now().isoformat(), description="创建时间，ISO格式的时间戳")
    metadata: Dict[str, Any] = Field(default_factory=dict, description="其他元数据，如词数、字符数等附加信息")

class TextCleaner:
    """文本清洗工具类
    
    负责对文本进行全面的清洗处理，包括去标识化、格式统一、单位标准化等操作，
    确保输出的文本符合后续处理的要求。
    """
    
    def __init__(self):
        # 用于去标识化的正则表达式字典
        # 定义了多种个人身份信息的匹配模式，用于后续的敏感信息替换
        self.patterns = {
            'id_card': r'[1-9]\d{5}(18|19|20)\d{2}(0[1-9]|1[0-2])(0[1-9]|[12]\d|3[01])\d{3}[0-9Xx]',  # 身份证号
            'phone': r'1[3-9]\d{9}|0\d{2,3}-\d{7,8}',  # 手机号和固定电话
            'name': r'[张王李赵刘陈杨黄周吴徐孙胡朱高林何郭马]{1,2}[\u4e00-\u9fa5]{1,2}',  # 常见中文姓名模式
            'date': r'\d{4}[-/]\d{1,2}[-/]\d{1,2}',  # 日期格式
            'hospital_id': r'医院编号[：:]?\s*\d{6,8}',  # 医院编号
            'medical_record': r'病历号[：:]?\s*\d{6,10}'  # 病历号
        }
        
    def remove_identifiers(self, text: str) -> str:
        """去除文本中的标识信息，保护患者隐私
        
        遍历预定义的正则表达式模式，将匹配到的敏感信息替换为标记文本，
        同时保留信息类型，便于后续分析和处理。
        
        参数:
            text: 需要处理的原始文本
        
        返回:
            去除标识信息后的文本
        """
        cleaned_text = text
        for pattern_name, pattern in self.patterns.items():
            # 将匹配到的敏感信息替换为大写的模式名称加上_REMOVED标记
            cleaned_text = re.sub(pattern, f'[{pattern_name.upper()}_REMOVED]', cleaned_text)
        return cleaned_text
    
    def normalize_format(self, text: str) -> str:
        """标准化文本格式：全半角转换、大小写统一、去除多余空白等
        
        将文本中的全角字符转换为半角字符，统一标点符号格式，
        并去除多余的空白字符，使文本格式更加规范统一。
        
        参数:
            text: 需要规范化的文本
        
        返回:
            格式规范化后的文本
        """
        # 全半角转换（简化版，将常见的全角标点符号转换为半角）
        text = text.replace('，', ',').replace('。', '.').replace('？', '?').replace('！', '!')
        text = text.replace('：', ':').replace('；', ';').replace('“', '"').replace('”', '"')
        text = text.replace('‘', "'").replace('’', "'").replace('（', '(').replace('）', ')')
        
        # 去除多余空白字符（包括空格、制表符、换行符等）
        text = re.sub(r'\s+', ' ', text).strip()
        
        return text
    
    def standardize_units(self, text: str) -> str:
        """标准化单位表示方式
        
        将文本中不同形式表示的单位统一为标准形式，便于后续的文本分析和数据提取。
        目前支持质量、体积、长度等常见医学单位的标准化。
        
        参数:
            text: 需要标准化单位的文本
        
        返回:
            单位标准化后的文本
        """
        # 单位映射字典，键为标准单位，值为需要替换的变体
        unit_mappings = {
            'kg': ['KG', 'Kg', '公斤'],  # 质量单位：千克
            'g': ['G', '克'],  # 质量单位：克
            'mg': ['MG', 'Mg', '毫克'],  # 质量单位：毫克
            'ml': ['ML', 'Ml', '毫升'],  # 体积单位：毫升
            'l': ['L', '升'],  # 体积单位：升
            'cm': ['CM', 'Cm', '厘米'],  # 长度单位：厘米
            'mm': ['MM', 'Mm', '毫米'],  # 长度单位：毫米
        }
        
        # 遍历映射字典，替换文本中的单位变体为标准单位
        for standard_unit, variations in unit_mappings.items():
            for variation in variations:
                text = text.replace(variation, standard_unit)
        
        return text
    
    def clean_text(self, text: str) -> str:
        """执行完整的文本清洗流程
        
        按顺序调用各个清洗方法，完成文本的全面清洗处理。
        清洗流程：去标识化 → 格式统一 → 单位标准化
        
        参数:
            text: 需要清洗的原始文本
        
        返回:
            经过完整清洗流程的文本
        """
        text = self.remove_identifiers(text)  # 第一步：去标识化
        text = self.normalize_format(text)    # 第二步：格式统一
        text = self.standardize_units(text)   # 第三步：单位标准化
        return text

class FileProcessor:
    """文件处理器，支持读取多种格式的文件
    
    负责读取PDF、DOCX、PPTX等不同格式的文件，提取文本内容，并将其转换为统一的
    文档片段格式，为后续的清洗和处理提供基础。
    """
    
    def __init__(self, cleaner: TextCleaner):
        """初始化文件处理器
        
        参数:
            cleaner: 文本清洗器实例，用于对提取的文本进行清洗处理
        """
        self.cleaner = cleaner
    
    def read_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        """读取PDF文件内容并提取文本片段
        
        使用pdfplumber库打开并读取PDF文件，逐页提取文本内容，
        对每一页的文本进行清洗，并创建对应的文档片段。
        
        参数:
            file_path: PDF文件的路径
        
        返回:
            文档片段列表，每个片段包含清洗后的文本和元数据
        """
        segments = []
        doc_id = str(uuid.uuid4())  # 为当前文档生成唯一标识符
        try:
            with pdfplumber.open(file_path) as pdf:
                # 逐页读取PDF内容
                for page_num, page in enumerate(pdf.pages, 1):
                    # 提取页面文本，如果无法提取则为空字符串
                    text = page.extract_text() or ""
                    # 跳过空文本或只包含空白字符的页面
                    if text.strip():
                        # 对提取的文本进行清洗
                        cleaned_text = self.cleaner.clean_text(text)
                        # 创建文档片段
                        segment = {
                            'doc_id': doc_id,
                            'segment_id': str(uuid.uuid4()),  # 为当前段落生成唯一标识符
                            'text': cleaned_text,
                            'original_path': file_path,
                            'file_type': 'pdf',
                            'page_num': page_num,
                            'section': f'Page {page_num}',
                            'metadata': {'word_count': len(cleaned_text.split())}  # 添加词数元数据
                        }
                        segments.append(segment)
            logger.info(f"成功读取PDF文件: {file_path}, 提取了 {len(segments)} 个段落")
        except Exception as e:
            # 记录读取过程中的错误
            logger.error(f"读取PDF文件失败: {file_path}, 错误: {str(e)}")
        return segments
    
    def read_docx(self, file_path: str) -> List[Dict[str, Any]]:
        """读取DOCX文件内容并提取文本片段
        
        使用python-docx库打开并读取DOCX文件，逐段提取文本内容，
        对每一段文本进行清洗，并创建对应的文档片段。
        
        参数:
            file_path: DOCX文件的路径
        
        返回:
            文档片段列表，每个片段包含清洗后的文本和元数据
        """
        segments = []
        doc_id = str(uuid.uuid4())  # 为当前文档生成唯一标识符
        try:
            doc = Document(file_path)
            # 逐段读取DOCX内容
            for para_idx, paragraph in enumerate(doc.paragraphs):
                text = paragraph.text
                # 跳过空文本或只包含空白字符的段落
                if text.strip():
                    # 对提取的文本进行清洗
                    cleaned_text = self.cleaner.clean_text(text)
                    # 创建文档片段
                    segment = {
                        'doc_id': doc_id,
                        'segment_id': str(uuid.uuid4()),  # 为当前段落生成唯一标识符
                        'text': cleaned_text,
                        'original_path': file_path,
                        'file_type': 'docx',
                        'section': f'Paragraph {para_idx+1}',
                        'metadata': {'word_count': len(cleaned_text.split())}  # 添加词数元数据
                    }
                    segments.append(segment)
            logger.info(f"成功读取DOCX文件: {file_path}, 提取了 {len(segments)} 个段落")
        except Exception as e:
            # 记录读取过程中的错误
            logger.error(f"读取DOCX文件失败: {file_path}, 错误: {str(e)}")
        return segments
    
    def read_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        """读取PPTX文件内容并提取文本片段
        
        使用python-pptx库打开并读取PPTX文件，逐页提取文本内容，
        对每一页的文本进行清洗，并创建对应的文档片段。
        
        参数:
            file_path: PPTX文件的路径
        
        返回:
            文档片段列表，每个片段包含清洗后的文本和元数据
        """
        segments = []
        doc_id = str(uuid.uuid4())  # 为当前文档生成唯一标识符
        try:
            prs = Presentation(file_path)
            # 逐页读取PPTX内容
            for slide_idx, slide in enumerate(prs.slides):
                slide_text = []
                # 提取幻灯片中所有形状的文本
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        slide_text.append(shape.text)
                # 将所有文本合并为一个字符串
                full_text = '\n'.join(slide_text)
                # 跳过空文本或只包含空白字符的幻灯片
                if full_text.strip():
                    # 对提取的文本进行清洗
                    cleaned_text = self.cleaner.clean_text(full_text)
                    # 创建文档片段
                    segment = {
                        'doc_id': doc_id,
                        'segment_id': str(uuid.uuid4()),  # 为当前段落生成唯一标识符
                        'text': cleaned_text,
                        'original_path': file_path,
                        'file_type': 'pptx',
                        'page_num': slide_idx + 1,
                        'section': f'Slide {slide_idx+1}',
                        'metadata': {'word_count': len(cleaned_text.split())}  # 添加词数元数据
                    }
                    segments.append(segment)
            logger.info(f"成功读取PPTX文件: {file_path}, 提取了 {len(segments)} 个段落")
        except Exception as e:
            # 记录读取过程中的错误
            logger.error(f"读取PPTX文件失败: {file_path}, 错误: {str(e)}")
        return segments
    
    def process_file(self, file_path: str) -> List[Dict[str, Any]]:
        """根据文件类型选择对应的处理器
        
        根据文件扩展名判断文件类型，调用相应的读取方法处理文件。
        目前支持PDF、DOCX、PPTX三种文件格式。
        
        参数:
            file_path: 要处理的文件路径
        
        返回:
            文档片段列表，每个片段包含清洗后的文本和元数据；
            如果文件类型不支持，则返回空列表
        """
        file_ext = os.path.splitext(file_path)[1].lower()  # 获取文件扩展名并转换为小写
        
        # 根据文件类型选择对应的读取方法
        if file_ext == '.pdf':
            return self.read_pdf(file_path)
        elif file_ext == '.docx':
            return self.read_docx(file_path)
        elif file_ext == '.pptx':
            return self.read_pptx(file_path)
        else:
            # 记录不支持的文件类型警告
            logger.warning(f"不支持的文件类型: {file_path}")
            return []

class DataIngestionPipeline:
    """数据摄入管道，协调整个处理流程
    
    负责协调整个数据处理流程，包括文件查找、文本提取、清洗、去重和结果保存等环节，
    是整个程序的核心控制类。
    """
    
    def __init__(self, input_dir: str, output_file: str):
        """初始化数据摄入管道
        
        参数:
            input_dir: 输入目录路径，包含要处理的源文件
            output_file: 输出文件路径，用于保存处理后的结果
        """
        self.input_dir = input_dir
        self.output_file = output_file
        self.cleaner = TextCleaner()
        self.processor = FileProcessor(self.cleaner)
        self.processed_hashes = set()  # 用于去重的文本哈希集合
    
    def find_files(self) -> List[str]:
        """查找输入目录下所有支持的文件类型
        
        遍历输入目录及其子目录，查找所有支持的文件格式（PDF、DOCX、PPTX），
        并返回文件路径列表。
        
        返回:
            支持格式的文件路径列表
        """
        supported_extensions = ('.pdf', '.docx', '.pptx')  # 支持的文件扩展名
        files = []
        
        # 遍历输入目录及其子目录
        for root, _, filenames in os.walk(self.input_dir):
            for filename in filenames:
                # 检查文件扩展名是否在支持的列表中
                if filename.lower().endswith(supported_extensions):
                    files.append(os.path.join(root, filename))  # 添加完整文件路径
        
        logger.info(f"找到 {len(files)} 个支持的文件")
        return files
    
    def is_duplicate(self, text: str) -> bool:
        """检查文本是否重复
        
        使用MD5哈希算法计算文本的哈希值，并与已处理文本的哈希集合进行比较，
        判断当前文本是否为重复内容。
        
        参数:
            text: 要检查的文本内容
        
        返回:
            如果文本重复则返回True，否则返回False
        """
        # 计算文本的MD5哈希值
        text_hash = hashlib.md5(text.encode('utf-8')).hexdigest()
        # 检查哈希值是否已存在于集合中
        if text_hash in self.processed_hashes:
            return True
        # 将新的哈希值添加到集合中
        self.processed_hashes.add(text_hash)
        return False
    
    def run(self):
        """运行整个数据摄入和清洗流程
        
        执行完整的数据处理流程：查找文件 → 处理每个文件 → 去重 → 保存结果。
        整个过程包括文本提取、清洗、验证和序列化等步骤。
        """
        logger.info(f"开始数据摄入和清洗流程，输入目录: {self.input_dir}")
        
        # 查找所有支持的文件
        files = self.find_files()
        
        # 处理每个文件
        all_segments = []
        for file_path in files:
            # 调用文件处理器处理当前文件
            segments = self.processor.process_file(file_path)
            # 去重：过滤掉重复的文本片段
            unique_segments = [seg for seg in segments if not self.is_duplicate(seg['text'])]
            # 将去重后的片段添加到总列表中
            all_segments.extend(unique_segments)
        
        # 保存结果到JSONL文件
        with open(self.output_file, 'w', encoding='utf-8') as f:
            for segment in all_segments:
                # 使用pydantic验证数据，确保数据格式正确
                doc_segment = DocumentSegment(**segment)
                # 将验证后的文档片段序列化为JSON字符串并写入文件
                f.write(doc_segment.model_dump_json() + '\n')
        
        logger.info(f"数据摄入和清洗流程完成，共处理 {len(all_segments)} 个段落，输出文件: {self.output_file}")

def main():
    """主函数，用于运行数据摄入和清洗流程
    
    负责程序的整体控制流程，包括参数配置、输入输出检查、用户交互和异常处理等。
    是程序的入口点，协调整个数据处理流程的执行。
    """
    # 配置参数
    INPUT_DIR = r"D:\rag-project\05-rag-practice\20-Data\04-临床文档"  # 输入目录路径
    OUTPUT_FILE = "clean_text.jsonl"  # 输出文件路径
    
    # 检查输入目录是否存在
    if not os.path.exists(INPUT_DIR):
        logger.error(f"输入目录不存在: {INPUT_DIR}")
        print(f"错误: 无法找到输入目录 '{INPUT_DIR}'")
        print("请检查目录路径是否正确，或者修改程序中的INPUT_DIR变量指向正确的临床文档目录")
        return
    
    # 检查目录中是否有支持的文件类型
    supported_extensions = ('.pdf', '.docx', '.pptx')
    found_files = []
    # 遍历输入目录查找支持的文件
    for root, _, filenames in os.walk(INPUT_DIR):
        for filename in filenames:
            if filename.lower().endswith(supported_extensions):
                found_files.append(os.path.join(root, filename))
    
    # 如果未找到支持的文件，给出警告并询问是否继续
    if not found_files:
        logger.warning(f"在目录 '{INPUT_DIR}' 中未找到支持的文件类型（PDF、DOCX、PPTX）")
        print(f"警告: 在目录 '{INPUT_DIR}' 中未找到支持的文件类型（PDF、DOCX、PPTX）")
        proceed = input("是否继续执行？(y/n): ")
        if proceed.lower() != 'y':
            return
    else:
        # 显示找到的文件数量
        print(f"在目录中找到 {len(found_files)} 个支持的文件")
    
    # 检查输出目录是否可写
    output_dir = os.path.dirname(os.path.abspath(OUTPUT_FILE))
    if not os.path.exists(output_dir):
        # 如果输出目录不存在，尝试创建
        try:
            os.makedirs(output_dir)
            logger.info(f"创建输出目录: {output_dir}")
        except Exception as e:
            # 记录创建目录失败的错误
            logger.error(f"创建输出目录失败: {str(e)}")
            print(f"错误: 无法创建输出目录 '{output_dir}': {str(e)}")
            return
    
    # 确认是否继续执行
    print(f"程序将从以下目录读取临床文档: {INPUT_DIR}")
    print(f"清洗后的文本将保存到: {OUTPUT_FILE}")
    print("开始处理...")
    
    # 创建并运行数据摄入管道
    try:
        # 初始化数据摄入管道
        pipeline = DataIngestionPipeline(INPUT_DIR, OUTPUT_FILE)
        # 运行数据处理流程
        pipeline.run()
        # 显示处理完成的消息
        print(f"处理完成！已成功生成清洗后的文本文件: {OUTPUT_FILE}")
    except Exception as e:
        # 捕获并记录执行过程中的错误
        logger.error(f"程序执行过程中发生错误: {str(e)}")
        print(f"错误: 程序执行过程中发生错误: {str(e)}")

if __name__ == "__main__":
    """程序入口点
    
    当作为主程序直接运行时，调用main()函数启动数据处理流程。
    """
    main()